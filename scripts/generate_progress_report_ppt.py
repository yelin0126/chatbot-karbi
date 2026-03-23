#!/usr/bin/env python3
"""
Generate a Korean PowerPoint progress report for the Tilon AI Chatbot project.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "reports"
OUTPUT_PATH = OUTPUT_DIR / "tilon_project_progress_20260323_ko.pptx"

BG = RGBColor(11, 18, 32)
PANEL = RGBColor(21, 32, 53)
ACCENT = RGBColor(26, 163, 154)
ACCENT_2 = RGBColor(245, 158, 11)
TEXT = RGBColor(241, 245, 249)
SUBTEXT = RGBColor(191, 205, 224)
MUTED = RGBColor(120, 144, 173)
GOOD = RGBColor(34, 197, 94)
WARN = RGBColor(245, 158, 11)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header_band(slide, title: str, subtitle: str | None = None) -> None:
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.33), Inches(0.8)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = PANEL
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(8.8), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(9.15), Inches(0.2), Inches(3.6), Inches(0.25))
        p = sub_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(11)
        r.font.color.rgb = SUBTEXT


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PANEL
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.6), Inches(10.8), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Tilon AI Chatbot 프로젝트 진행 보고"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = TEXT

    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "문서 기반 RAG 고도화, 업로드 워크플로우 안정화, QLoRA 준비 현황"
    r.font.size = Pt(16)
    r.font.color.rgb = SUBTEXT

    info_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(5.2), Inches(1.6))
    tf = info_box.text_frame
    for idx, line in enumerate(
        [
            "보고일: 2026-03-23",
            "대상: 팀 리더 / 내부 진행 공유",
            "프로젝트 상태: 중간 고도화 단계",
        ]
    ):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(18 if idx == 0 else 16)
        r.font.color.rgb = TEXT if idx == 0 else SUBTEXT

    callout = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(3.0), Inches(5.2), Inches(1.6)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(16, 28, 46)
    callout.line.color.rgb = ACCENT
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "핵심 요약\nRAG 기반 문서 QA는 안정권에 진입했고,\n현재 주력 과제는 성능/평가 고도화와 QLoRA 학습 준비입니다."
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = TEXT


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, title, subtitle)

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(12.0), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, title)

    for x, col_title, items, border in [
        (0.65, left_title, left_items, ACCENT),
        (6.9, right_title, right_items, ACCENT_2),
    ]:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.25), Inches(5.7), Inches(5.8)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL
        panel.line.color.rgb = border

        title_box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(1.45), Inches(5.1), Inches(0.4))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = col_title
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = TEXT

        body = slide.shapes.add_textbox(Inches(x + 0.28), Inches(1.95), Inches(5.0), Inches(4.7))
        tf = body.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = SUBTEXT
            p.space_after = Pt(8)


def add_metrics_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "평가 및 수치 현황", "2026-03-23 기준")

    metrics = [
        ("라이브러리 벤치마크", "안정화", "기존 문서 QA 기준선 확보"),
        ("업로드 벤치마크", "8 rows / 소스 재현성 1.0", "번들 PDF 모호성/거절 케이스 포함"),
        ("업로드 평균 답변 recall", "0.896", "핵심 포인트 회수율 양호"),
        ("검색 지연시간", "Hybrid 8.9ms", "업로드 기준"),
        ("재정렬 지연시간", "Hybrid+Rerank 6692ms", "정확도 향상 있으나 실시간 채팅에는 과도"),
        ("QLoRA 학습 데이터", "25 rows", "한글 20 / 영어 5"),
    ]

    top = Inches(1.4)
    left = Inches(0.8)
    for idx, (label, value, note) in enumerate(metrics):
        row = idx // 2
        col = idx % 2
        x = left + Inches(col * 6.1)
        y = top + Inches(row * 1.75)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(5.55), Inches(1.35)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = ACCENT if idx % 2 == 0 else ACCENT_2

        label_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12), Inches(5.1), Inches(0.28))
        p = label_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.size = Pt(14)
        r.font.color.rgb = MUTED

        value_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.42), Inches(5.1), Inches(0.36))
        p = value_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = value
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = TEXT

        note_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.86), Inches(5.1), Inches(0.28))
        p = note_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = note
        r.font.size = Pt(11)
        r.font.color.rgb = SUBTEXT


def add_status_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "현재 단계와 병목")

    statuses = [
        ("완료/강점", GOOD, [
            "문서 업로드/파싱/청킹/벡터저장 구조 안정화",
            "하이브리드 검색 + 문서 스코프 QA 동작",
            "업로드 재선택 UI, 다중 업로드, 비교용 문서 선택 지원",
            "라이브러리/업로드 벤치마크 기반선 확보",
        ]),
        ("현재 병목", WARN, [
            "Layer 9B: 검색 지연시간/런타임 최적화",
            "재정렬기는 품질 향상은 있으나 실시간 비용이 큼",
            "비교형 PDF QA는 구현되었지만 검증은 시작 단계",
            "QLoRA는 학습 스크립트/데이터 준비 단계",
        ]),
    ]

    for idx, (title, color, lines) in enumerate(statuses):
        x = Inches(0.8 + idx * 6.2)
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(5.6), Inches(4.8)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL
        panel.line.color.rgb = color

        title_box = slide.shapes.add_textbox(x + Inches(0.22), Inches(1.7), Inches(5.0), Inches(0.4))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = TEXT

        body = slide.shapes.add_textbox(x + Inches(0.25), Inches(2.2), Inches(5.0), Inches(3.6))
        tf = body.text_frame
        for j, line in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = SUBTEXT
            p.space_after = Pt(8)


def add_timeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "다음 단계 제안")

    steps = [
        ("1", "검색/비교 안정화", "비교형 PDF 벤치마크 실행\n비교 응답 품질 점검 및 보정"),
        ("2", "평가셋 확대", "한글 중심으로 질문셋 확대\n영문/혼합 질의는 보조적으로 추가"),
        ("3", "QLoRA 준비 강화", "학습 데이터 25 -> 40+ 확대\n강한 정답/거절/비교 케이스 보강"),
        ("4", "첫 학습 실행", "기준 모델 고정 후\n첫 QLoRA adapter 학습 및 비교"),
    ]

    for idx, (num, title, detail) in enumerate(steps):
        x = Inches(0.8 + idx * 3.1)
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.95), Inches(1.35), Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if idx % 2 == 0 else ACCENT_2
        circle.line.fill.background()
        p = circle.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = BG

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(2.65), Inches(3.15)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL
        box.line.color.rgb = ACCENT if idx % 2 == 0 else ACCENT_2

        title_box = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.52), Inches(2.35), Inches(0.4))
        p = title_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = TEXT

        detail_box = slide.shapes.add_textbox(x + Inches(0.16), Inches(3.0), Inches(2.3), Inches(1.9))
        p = detail_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = detail
        r.font.size = Pt(14)
        r.font.color.rgb = SUBTEXT


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "결론")

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.6), Inches(11.5), Inches(4.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = ACCENT

    body = slide.shapes.add_textbox(Inches(1.25), Inches(2.0), Inches(10.8), Inches(3.9))
    tf = body.text_frame
    lines = [
        "현재 프로젝트는 초기 프로토타입 단계를 넘어서, 문서 기반 RAG 제품의 핵심 구조를 확보한 상태입니다.",
        "특히 업로드 문서 처리, 문서 스코프 QA, 벤치마크 기반 평가, 비교형 질문 지원까지 연결되었습니다.",
        "다음 핵심 과제는 한글 중심 평가셋 확대와 QLoRA 학습 루프 시작입니다.",
        "즉, 지금은 기능 추가보다 '품질 측정 -> 데이터 확장 -> 미세조정'으로 넘어가야 하는 시점입니다.",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(22 if idx == 0 else 18)
        p.font.color.rgb = TEXT if idx == 0 else SUBTEXT
        p.space_after = Pt(12)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_bullets_slide(
        prs,
        "프로젝트 개요",
        [
            "목표: 한국어/영어 PDF 및 이미지 기반 문서 QA 챗봇 구축",
            "핵심 구성: 파싱/OCR -> 청킹/임베딩 -> 하이브리드 검색 -> 근거 기반 답변",
            "확장 방향: 업로드 문서 QA, 다중 문서 비교, QLoRA 기반 응답 품질 향상",
        ],
    )
    add_two_column_slide(
        prs,
        "현재까지 구현된 핵심 기능",
        "완료된 기능",
        [
            "텍스트 PDF / 스캔 PDF / 이미지 업로드 처리",
            "업로드 문서 기억 및 사이드바 재선택",
            "동일 파일 재업로드 시 중복 청크 정리",
            "번들 PDF 모호성 감지 및 서브지침 범위 축소",
        ],
        "최근 확장 기능",
        [
            "다중 업로드 및 선택 기반 비교형 질문 지원",
            "비교형 벤치마크 파일 추가",
            "QLoRA 학습용 초기 데이터셋 확장",
            "실시간 지연시간을 고려한 하이브리드 검색 정책",
        ],
    )
    add_metrics_slide(prs)
    add_status_slide(prs)
    add_bullets_slide(
        prs,
        "UI 및 사용자 흐름 현황",
        [
            "웹 UI 기준으로 한 번에 여러 파일 업로드 가능",
            "업로드한 문서는 좌측 사이드바에 기억되어 재사용 가능",
            "여러 문서를 선택하면 비교 질문 흐름으로 전환",
            "현재 UI는 내부 테스트용으로 충분하나, 운영용 완성도는 추가 개선 여지 존재",
        ],
    )
    add_bullets_slide(
        prs,
        "QLoRA 준비 현황",
        [
            "학습 스크립트와 의존성 파일은 준비 완료",
            "학습 데이터셋은 25개 샘플까지 확장 완료",
            "현재 비중은 한국어 중심(20) + 영어 보조(5)",
            "다음 단계는 비교/거절/영문 케이스를 조금 더 늘린 뒤 첫 adapter 학습 실행",
        ],
    )
    add_timeline_slide(prs)
    add_closing_slide(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(path)
